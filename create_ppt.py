from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

TEMPLATE = r'c:\Users\qlg9915\AnypointStudio\ws-5\amq-handler-poc\Takeda.pptx'
prs = Presentation(TEMPLATE)

# Note: Template slides remain at the end. Delete them manually in PowerPoint after opening.

# Colors
DARK = RGBColor(0x1A, 0x1A, 0x2E)
BLUE = RGBColor(0x00, 0xA0, 0xDF)
ACCENT = RGBColor(0x4A, 0x6C, 0xF7)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Layout references
LY_TITLE = prs.slide_layouts[0]      # Title Page 1
LY_CONTENT = prs.slide_layouts[7]    # Standard 1-Column Text
LY_DIVIDER = prs.slide_layouts[6]    # 3_Divider 1
LY_2COL = prs.slide_layouts[8]       # 1_Standard 1-Column Text (2 col)
LY_END = prs.slide_layouts[9]        # End Slide

def set_ph(slide, idx, text, font_size=None, bold=None, color=None):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            if font_size or bold is not None or color:
                for p in ph.text_frame.paragraphs:
                    for r in p.runs:
                        if font_size: r.font.size = Pt(font_size)
                        if bold is not None: r.font.bold = bold
                        if color: r.font.color.rgb = color
            return ph
    return None

def set_ph_bullets(slide, idx, items, font_size=14, color=DARK):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.clear()
            for i, item in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)
                p.font.color.rgb = color
                p.space_after = Pt(6)
            return ph
    return None

# ===================== SLIDE 1: Title =====================
slide = prs.slides.add_slide(LY_TITLE)
set_ph(slide, 0, 'AMQ Studio', font_size=40, bold=True)
set_ph(slide, 27, 'A Self-Service Queue Management Dashboard for Anypoint MQ', font_size=16)
set_ph(slide, 26, 'Proof of Concept  |  Inspired by Anypoint MQ REM  |  Extended with Advanced Features', font_size=12, color=GRAY)
set_ph(slide, 23, 'Built with MuleSoft  |  CloudHub Ready  |  Feedback Welcome', font_size=10, color=GRAY)

# ===================== SLIDE 2: Divider - Background =====================
slide = prs.slides.add_slide(LY_DIVIDER)
set_ph(slide, 0, 'Background')
set_ph(slide, 27, 'Error handling in message-driven integrations')

# ===================== SLIDE 3: Background - Error Handling Pattern =====================
slide = prs.slides.add_slide(LY_CONTENT)
set_ph(slide, 0, 'The Error Handling Pattern in Anypoint MQ')
background = [
    'Message-driven integrations follow a common error handling pattern:',
    '',
    '  1. An error occurs during message processing — caught by the error handler',
    '  2. The original message is enriched with retry metadata (delivery count, error info)',
    '  3. The enriched message is published to a dedicated error queue',
    '  4. After a delay, messages in the error queue are reprocessed:',
    '       • If retry count < limit → re-queued to the source queue for another attempt',
    '       • If retry limit exceeded → moved to a Dead Letter Queue (DLQ) for investigation',
    '',
    'The challenge: Once messages land in the DLQ, MuleSoft does not have native',
    'functionality to manage, inspect, or resubmit them. Teams are left writing',
    'one-off scripts or manually handling failed messages.',
]
set_ph_bullets(slide, 15, background, font_size=13)

# ===================== SLIDE 4: Divider - Problem =====================
slide = prs.slides.add_slide(LY_DIVIDER)
set_ph(slide, 0, 'The Problem')
set_ph(slide, 27, 'Gaps in Anypoint MQ message management')

# ===================== SLIDE 5: Problem Statement =====================
slide = prs.slides.add_slide(LY_CONTENT)
set_ph(slide, 0, 'Problem Statement')
problems = [
    'MuleSoft does not have native functionality to manage messages on a queue',
    'in terms of inspection, resubmission, or transfer to other queues.',
    '',
    '❌  No visibility — Cannot view message payloads without consuming (and losing) them',
    '❌  No resubmission — No built-in way to return messages from DLQ to source queue',
    '❌  No cross-queue transfer — Moving messages between queues requires custom code',
    '❌  No editing — Corrupted payloads cannot be fixed before reprocessing',
    '❌  No multi-format awareness — JSON, CSV, XML, text all treated the same',
    '❌  Manual effort — Teams write throwaway scripts, wasting time and risking message loss',
    '❌  No concurrency safety — Multiple users accessing the same queue can cause conflicts',
]
set_ph_bullets(slide, 15, problems, font_size=13)

# ===================== SLIDE 6: Inspiration =====================
slide = prs.slides.add_slide(LY_CONTENT)
set_ph(slide, 0, 'Inspiration: Anypoint MQ REM')
inspiration = [
    'MuleSoft\'s Anypoint MQ REM (Resubmit Error Messages) introduced a dashboard',
    'concept — a Web API built with MuleSoft Platform APIs and Anypoint Studio that:',
    '',
    '  ✅  Views messages in the error/dead-letter queue',
    '  ✅  Resubmits individual messages or all messages to the source queue',
    '  ✅  Refreshes to see new messages in real time',
    '  ✅  Shows statistics — message counts, timestamps, retry frequency',
    '',
    'AMQ Studio takes this concept further by adding:',
    '',
    '  🆕  Inline payload editing (JSON, CSV table editor, XML, text)',
    '  🆕  Cross-queue transfers (not just return-to-origin)',
    '  🆕  Message deletion with double confirmation',
    '  🆕  Concurrency locking with auto-return on idle',
    '  🆕  Multi-format content type preservation through round-trips',
    '  🆕  Dynamic Platform API integration (cascading org/env/region/queue)',
    '  🆕  AES-256 encrypted session security',
]
set_ph_bullets(slide, 15, inspiration, font_size=12)

# ===================== SLIDE 7: Divider - Solution =====================
slide = prs.slides.add_slide(LY_DIVIDER)
set_ph(slide, 0, 'The Solution')
set_ph(slide, 27, 'AMQ Studio — extending REM with advanced capabilities')

# ===================== SLIDE 8: Solution Overview =====================
slide = prs.slides.add_slide(LY_CONTENT)
set_ph(slide, 0, 'Solution: AMQ Studio (Proof of Concept)')
solution = [
    'A self-contained Mule application with a built-in web dashboard for Anypoint MQ',
    'queue management — deployable to CloudHub, Runtime Fabric, or on-premise.',
    '',
    '🖥️  FRONTEND — HTML/JS/CSS served via parse-template, AES-256 encrypted sessions',
    '⚙️  BACKEND — 7 modular Mule XML files, Object Store for temporary message storage',
    '☁️  PLATFORM — Proxy to Anypoint Platform APIs for dynamic configuration',
    '',
    'TWO ACCESS MODES (PoC):',
    '  • Platform-led login — Authenticate with Connected App, then Org/Env/Region/Queue',
    '    auto-populate via Platform APIs (cascading dropdowns)',
    '  • Manual form — Enter details directly for quick access (will be removed at GA)',
    '',
    'USER FLOW:',
    '  1. Authenticate with Connected App credentials (AES-256-GCM encrypted)',
    '  2. Select Organization → Environment → Region → Queue (auto-populated)',
    '  3. Lock acquired → Messages consumed → Displayed progressively in real-time',
    '  4. Inspect, edit, transfer, delete, or return messages',
    '  5. On disconnect → All messages safely returned → Lock released → Session cleared',
]
set_ph_bullets(slide, 15, solution, font_size=12)

# ===================== SLIDE 9: Use Cases =====================
slide = prs.slides.add_slide(LY_CONTENT)
set_ph(slide, 0, 'Use Cases')
use_cases = [
    '🔍  DLQ Inspection — View all messages in error/dead-letter queues with full payload',
    '🔄  Message Resubmission — Return individual or all messages to the source queue',
    '↗️  Cross-Queue Transfer — Move messages to a different queue (same region/org/env)',
    '✏️  Payload Editing — Fix corrupted JSON, edit CSV in a table view, modify properties',
    '🗑️  Message Cleanup — Delete obsolete messages with double confirmation safety',
    '📊  Problem Discovery — Identify which messages fail most (via payload inspection)',
    '⚡  Real-time Monitoring — Progressive streaming shows messages as they arrive',
    '🔒  Safe Multi-User Access — Lock mechanism prevents concurrent queue access',
]
set_ph_bullets(slide, 15, use_cases, font_size=14)

# ===================== SLIDE 10: Architecture =====================
slide = prs.slides.add_slide(LY_CONTENT)
set_ph(slide, 0, 'Architecture & File Structure')
arch = [
    'amq-handler-poc.xml  →  Global configs (connectors, object stores, HTTP listener)',
    'common.xml  →  Shared sub-flows (OS store, retrieve, delete)',
    'message-operations.xml  →  Fetch, list, view, update, delete messages',
    'queue-operations.xml  →  Return to origin, transfer to another queue',
    'lock-management.xml  →  Acquire, renew, release lock + stale lock scheduler',
    'platform-proxy.xml  →  Anypoint Platform API proxies (orgs, envs, regions, queues)',
    'static-assets.xml  →  Serve HTML, CSS, JS, SVG via parse-template',
    '',
    'CONFIGURATION (dev.yaml):',
    '  • All values configurable — URLs, timeouts, TTLs, scheduler frequency',
    '  • Organization-agnostic — deploy to any org, works with any credentials',
    '  • Test mode toggle via property: app.testModeEnabled',
]
set_ph_bullets(slide, 15, arch, font_size=13)

# ===================== SLIDE 11: Security =====================
slide = prs.slides.add_slide(LY_CONTENT)
set_ph(slide, 0, 'Security Features')
security = [
    '🔐  CLIENT-SIDE SECURITY',
    '  • AES-256-GCM encryption for stored credentials (non-extractable keys)',
    '  • sessionStorage — cleared on tab close, secrets never pre-filled',
    '  • Credentials never stored in plaintext anywhere',
    '',
    '🛡️  SERVER-SIDE SECURITY',
    '  • Queue locking prevents concurrent access by multiple users',
    '  • Auto-return on idle prevents data loss (configurable threshold)',
    '  • Object Store TTL ensures automatic cleanup',
    '  • Session cleared on disconnect — nothing persists',
    '',
    '✅  OPERATIONAL SAFETY',
    '  • Double confirmation for destructive actions (type DELETE / SAVE to proceed)',
    '  • Session termination returns all untouched messages to queue',
    '  • No data loss scenario — messages always safe',
    '',
    '📦  DEPLOYMENT',
    '  • JS/CSS minified via Terser + CleanCSS in Maven build',
    '  • CloudHub / Runtime Fabric / On-prem ready | LTS compatible (Mule 4.9.x)',
]
set_ph_bullets(slide, 15, security, font_size=12)

# ===================== SLIDE 12: REM vs AMQ Studio =====================
slide = prs.slides.add_slide(LY_CONTENT)
set_ph(slide, 0, 'REM vs AMQ Studio — Feature Comparison')
comparison = [
    'FEATURE                              REM          AMQ STUDIO',
    '─────────────────────────────────────────────────────────────',
    'View DLQ messages                      ✅               ✅',
    'Resubmit to source queue               ✅               ✅',
    'Refresh for new messages               ✅               ✅',
    'Message statistics                      ✅               ✅',
    'Transfer to different queue             ❌               ✅',
    'Edit message payload                    ❌               ✅',
    'Inline CSV table editing                ❌               ✅',
    'Delete messages                         ❌               ✅',
    'Multi-format preservation               ❌               ✅',
    'Concurrency locking                     ❌               ✅',
    'Auto-return on idle                     ❌               ✅',
    'Real-time progressive streaming        ❌               ✅',
    'AES-256 encrypted sessions             ❌               ✅',
    'Platform API cascading dropdowns       ❌               ✅',
    'Minified/obfuscated distribution       ❌               ✅',
]
set_ph_bullets(slide, 15, comparison, font_size=11)

# ===================== SLIDE 13: Divider - Demo =====================
slide = prs.slides.add_slide(LY_DIVIDER)
set_ph(slide, 0, 'Live Demo')
set_ph(slide, 27, 'http://localhost:8081/amq-studio')

# ===================== SLIDE 14: Roadmap =====================
slide = prs.slides.add_slide(LY_CONTENT)
set_ph(slide, 0, 'Roadmap & Future Enhancements')
roadmap = [
    '🔐  SSO Login — Connected App OAuth 2.0 flow, "Login with MuleSoft" button',
    '📊  Analytics Dashboard — Message volume trends, queue health monitoring',
    '🔄  Exchange & DLQ — Topic/exchange support, DLQ-specific workflows, bulk replay',
    '📈  Error Intelligence — Track retry patterns, identify recurring failures',
    '🧪  Testing & CI/CD — MUnit coverage, GitHub Actions, automated CloudHub deploy',
    '📱  UX — Dark mode, mobile responsive, keyboard shortcuts',
    '🔌  Extensibility — Plugin architecture for custom actions on messages',
]
set_ph_bullets(slide, 15, roadmap, font_size=14)

# ===================== SLIDE 15: Feedback Request =====================
slide = prs.slides.add_slide(LY_CONTENT)
set_ph(slide, 0, 'We Need Your Feedback!')
feedback = [
    'This is a Proof of Concept — your feedback will shape the final product.',
    '',
    '💬  What works well?',
    '🔧  What could be improved?',
    '🆕  What features would you like to see next?',
    '🐛  Did you encounter any issues?',
    '',
    'CURRENT STATUS:',
    '  • Platform-led login — fully functional (auto-populates via Platform APIs)',
    '  • Manual form — available for quick access during PoC (will be removed at GA)',
    '  • Security — AES-256-GCM encryption, auto-return, double confirmations',
    '',
    'Please try it out and share your thoughts!',
    'Every piece of feedback helps shape the final release.',
]
set_ph_bullets(slide, 15, feedback, font_size=13)

# ===================== SLIDE 16: End =====================
slide = prs.slides.add_slide(LY_END)
set_ph(slide, 10, 'Thank You!\n\nAMQ Studio — Self-Service Queue Management for Anypoint MQ\nProof of Concept  |  Inspired by Anypoint MQ REM\n\n🤝 Your Feedback Makes This Better\n\nQuestions?', font_size=22)

# Save
output_path = r'c:\Users\qlg9915\AnypointStudio\ws-5\amq-handler-poc\AMQ_Studio_Demo.pptx'
prs.save(output_path)
print(f'Saved to: {output_path}')

